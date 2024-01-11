from PIL import Image, ImageOps
import io
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

def invert_image_colors(slide):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Extract the image from the shape
            image_stream = io.BytesIO(shape.image.blob)
            with Image.open(image_stream) as img:
                # If image format is not supported convert it first
                if img.mode not in ('RGB', 'RGBA'):
                    img = img.convert('RGB')

                # Invert the image color using Pillow only if supported
                inverted_img = ImageOps.invert(img)

                # Save the inverted image to a stream
                img_stream = io.BytesIO()
                inverted_img.save(img_stream, format='PNG')
                img_stream.seek(0)

                # Remove the original picture
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                slide.shapes._spTree.remove(shape._element)

                # Add the new inverted image to the slide
                pic = slide.shapes.add_picture(img_stream, left, top, width, height)

