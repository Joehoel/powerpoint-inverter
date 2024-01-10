from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

def invert_text_and_background(slide):
    # Invert slide background to black
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    # Invert text color to white for all shapes that contain text
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
